from __future__ import annotations

import shutil
from pathlib import Path

import pymupdf
from pptx import Presentation

from presentation_video.domain.models import PresentationDocument, SlideContent
from presentation_video.domain.ports import DocumentIngestor
from presentation_video.infrastructure.process import run_process


def _first_non_empty_line(text: str) -> str:
    return next((line.strip() for line in text.splitlines() if line.strip()), "")


def _render_pdf(pdf_path: Path, output_dir: Path, dpi: int = 160) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    rendered: list[Path] = []
    document = pymupdf.open(pdf_path)
    try:
        for index in range(1, document.page_count + 1):
            page = document.load_page(index - 1)
            output_path = output_dir / f"slide-{index:03d}.png"
            page.get_pixmap(dpi=dpi, alpha=False).save(output_path)
            rendered.append(output_path)
    finally:
        document.close()
    return rendered


class PdfIngestor(DocumentIngestor):
    async def ingest(self, source: Path, work_dir: Path) -> PresentationDocument:
        images = _render_pdf(source, work_dir / "slides")
        document = pymupdf.open(source)
        try:
            page_features: list[tuple[str, int, int]] = []
            for page_index in range(document.page_count):
                page = document.load_page(page_index)
                text = page.get_text("text").strip()
                page_features.append(
                    (text, len(page.get_images(full=True)), len(page.get_drawings()))
                )
            text_heavy_pages = sum(
                len(text.split()) >= 120 and drawing_count < 3
                for text, _, drawing_count in page_features
            )
            document_is_long_form_text = (
                bool(page_features)
                and text_heavy_pages / len(page_features) >= 0.65
            )
            slides: list[SlideContent] = []
            for index in range(1, document.page_count + 1):
                text, embedded_images, vector_drawings = page_features[index - 1]
                word_count = len(text.split())
                source_frame_suitable = not (
                    document_is_long_form_text
                    or (
                        word_count >= 120
                        and vector_drawings < 3
                        and embedded_images <= 1
                    )
                )
                slides.append(
                    SlideContent(
                        number=index,
                        title=_first_non_empty_line(text),
                        body_text=text,
                        image_path=images[index - 1],
                        source_frame_suitable=source_frame_suitable,
                    )
                )
        finally:
            document.close()
        return PresentationDocument(
            source_path=source,
            title=slides[0].title if slides else source.stem,
            slides=slides,
        )


class PptxIngestor(DocumentIngestor):
    async def ingest(self, source: Path, work_dir: Path) -> PresentationDocument:
        libreoffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not libreoffice:
            raise RuntimeError(
                "LibreOffice is required to render PPTX files. Install it and expose 'soffice' in PATH."
            )

        converted_dir = work_dir / "converted"
        converted_dir.mkdir(parents=True, exist_ok=True)
        await run_process(
            libreoffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(converted_dir),
            str(source),
        )
        pdf_path = converted_dir / f"{source.stem}.pdf"
        if not pdf_path.exists():
            raise RuntimeError(f"LibreOffice did not create the expected PDF: {pdf_path}")

        images = _render_pdf(pdf_path, work_dir / "slides")
        presentation = Presentation(str(source))
        slides: list[SlideContent] = []

        for index, slide in enumerate(presentation.slides, start=1):
            text_parts: list[str] = []
            title = ""
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        text_parts.append(text)
                        if not title:
                            title = _first_non_empty_line(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(" |"):
                            text_parts.append(row_text)

            notes = ""
            try:
                notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
                if notes_frame is not None:
                    notes = notes_frame.text.strip()
            except (AttributeError, ValueError):
                notes = ""

            slides.append(
                SlideContent(
                    number=index,
                    title=title or f"Slide {index}",
                    body_text="\n".join(text_parts),
                    speaker_notes=notes,
                    image_path=images[index - 1],
                    source_frame_suitable=True,
                )
            )

        return PresentationDocument(
            source_path=source,
            title=slides[0].title if slides else source.stem,
            slides=slides,
        )


class ExtensionDocumentIngestorFactory:
    def __init__(self) -> None:
        self._strategies: dict[str, DocumentIngestor] = {
            ".pdf": PdfIngestor(),
            ".pptx": PptxIngestor(),
        }

    def create(self, source: Path) -> DocumentIngestor:
        extension = source.suffix.lower()
        try:
            return self._strategies[extension]
        except KeyError as exc:
            supported = ", ".join(sorted(self._strategies))
            raise ValueError(f"Unsupported file type '{extension}'. Supported: {supported}") from exc
